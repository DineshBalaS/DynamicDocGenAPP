import re
from io import BytesIO
from copy import deepcopy
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.util import Inches

def _transfer_font_properties(source_font, target_font):
    """
    A helper function to manually copy key font properties from one
    run to another. This is a robust way to preserve styling.
    """
    target_font.name = source_font.name
    target_font.size = source_font.size
    target_font.bold = source_font.bold
    target_font.italic = source_font.italic
    target_font.underline = source_font.underline
    # Handle all common color types, not just RGB.
    
    # Get the color object from the source font
    source_color = source_font.color
    
    # Case 1: RGB Color (e.g., "Red" from standard colors)
    if source_color.type == MSO_COLOR_TYPE.RGB:
        target_font.color.rgb = source_color.rgb
        
    # Case 2: Theme Color (e.g., "White, Background 1" or "Blue, Accent 1")
    elif source_color.type == MSO_COLOR_TYPE.SCHEME:
        target_font.color.theme_color = source_color.theme_color
        # Also copy brightness adjustments (tint/shade)
        target_font.color.brightness = source_color.brightness
        
    # Case 3: Preset Color (less common, but good to handle)
    elif source_color.type == MSO_COLOR_TYPE.PRESET:
        target_font.color.preset_color = source_color.preset_color

def extract_placeholders(file_stream: BytesIO) -> list:
    """
    [cite_start]Parses a .pptx file stream to find unique placeholders[cite: 296].

    This function is used by the /api/upload endpoint to analyze a template
    [cite_start]without needing to search inside complex elements like tables or charts[cite: 299, 300].

    Args:
        [cite_start]file_stream: A file-like object representing the .pptx file[cite: 296].

    Returns:
        A list of unique dictionaries, where each dictionary represents a
        [cite_start]placeholder with its "name" and "type"[cite: 305, 306].
        Example: [{"name": "client", "type": "text"}, {"name": "logo", "type": "image"}]
    """
    # [cite_start]Regex to find placeholders in two formats: {{name}} and {{type:name}}[cite: 302].
    pattern = re.compile(r'\{\{(?:(\w+):)?(\w+)\}\}')
    
    # [cite_start]Use a set to store unique (name, type) tuples to handle duplicates[cite: 307].
    found_placeholders = set()

    try:
        prs = Presentation(file_stream)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                for paragraph in shape.text_frame.paragraphs:
                    # Find all matches in the paragraph's text
                    matches = pattern.findall(paragraph.text)
                    for ph_type, ph_name in matches:
                        # [cite_start]Default to "text" if the type is not specified[cite: 304].
                        final_type = ph_type if ph_type else "text"
                        found_placeholders.add((ph_name, final_type))

    except Exception as e:
        print(f"Error extracting placeholders: {e}")
        # In a real app, you might raise a custom exception here
        raise ValueError("Could not process the presentation file.")

    # Convert the set of tuples to a sorted list of dictionaries for consistent output.
    return sorted(
        [{"name": name, "type": type_} for name, type_ in found_placeholders],
        key=lambda x: x['name']
    )
    
def generate_presentation(template_stream: BytesIO, data: dict, s3_service) -> BytesIO:
    """
    Generates a presentation by manually replacing placeholders in a template stream.
    This function uses the base python-pptx library for all manipulations.
    """
    ppt = Presentation(template_stream)
    # Regex to find list placeholders specifically
    list_pattern = re.compile(r'\{\{list:(\w+)\}\}')
    # Regex to find image/scrape placeholders specifically
    image_pattern = re.compile(r'\{\{(?:image|scrape):(\w+)\}\}')
    # Regex for simple text placeholders (including explicitly typed text ones)
    text_pattern = re.compile(r'\{\{(?:text:|choice:)?(\w+)\}\}')

    for slide in ppt.slides:
        shapes_to_delete = []
        for shape in list(slide.shapes): # Use list() to allow safe deletion
            if not shape.has_text_frame:
                continue
            
            # --- Check shape text content ONCE ---
            shape_text = shape.text_frame.text

            # --- Image Replacement Logic ---
            if '{{image:' in shape_text or '{{scrape:' in shape_text:
                match = image_pattern.search(shape_text)
                if match:
                    ph_name = match.group(1)
                    s3_key = data.get(ph_name)
                    
                    if s3_key:
                        try:
                            image_stream = s3_service.download_file_as_stream(s3_key)
                            slide.shapes.add_picture(
                                image_stream, shape.left, shape.top, 
                                width=shape.width, height=shape.height
                            )
                            shapes_to_delete.append(shape)
                        except Exception as e:
                            print(f"ERROR: Could not add image for '{ph_name}'. Details: {e}")
                            shape.text_frame.text = f"[Image Error]"
                    else:
                        shape.text_frame.text = f"[Image Missing: {ph_name}]"
                    continue # Skip other replacements for this shape

            # --- List Replacement Logic ---
            found_list_in_shape = False
            target_para_idx = -1
            source_font = None
            list_ph_name = None
            target_para_obj = None # Store the paragraph object itself

            # Find the paragraph containing the list placeholder first
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                list_match = list_pattern.search(para.text)
                if list_match:
                    if para.runs:
                        source_font = para.runs[0].font # Store font from the first run
                    else:
                        source_font = None
                    target_para_idx = para_idx
                    target_para_obj = para # Keep the paragraph object
                    list_ph_name = list_match.group(1)
                    break # Found the paragraph, stop searching

            # Process the list if a placeholder paragraph was found
            if target_para_idx != -1 and list_ph_name is not None and target_para_obj is not None:
                items = data.get(list_ph_name, []) # Expect data[list_ph_name] to be a list
                tf = shape.text_frame # Get the text frame
                
                # We will re-use the original placeholder paragraph for the first item.
                # This preserves all paragraph formatting (bullets, indentation, etc.).
                p = target_para_obj 
                
                if items and isinstance(items, list) and any(str(item).strip() for item in items):
                    # Filter out any empty strings
                    valid_items = [str(item) for item in items if str(item).strip()]
                    
                    # 1. Set the text for the first item (in the existing paragraph)
                    p.clear() # Clear existing runs (like '{{list:name}}')
                    run = p.add_run()
                    run.text = valid_items[0]
                    if source_font:
                        _transfer_font_properties(source_font, run.font)
                    
                    # 2. Add subsequent items as new paragraphs
                    for item_text in valid_items[1:]:
                        # Add a new paragraph element
                        new_p = tf.add_paragraph()

                       
                        # Get the <a:pPr> (paragraph properties) element from the original paragraph
                        pPr_to_copy = p._element.pPr

                        # Only proceed if the original paragraph *has* properties to copy
                        if pPr_to_copy is not None:
                            # Get the <a:pPr> element of the new paragraph,
                            # CREATING IT if it doesn't exist. This is the fix.
                            new_pPr = new_p._element.get_or_add_pPr()
                            
                            # Clear any default properties that might be on the new pPr
                            new_pPr.clear()
                            
                            # Copy all XML attributes (like 'lvl', 'marL', etc.)
                            new_pPr.attrib.update(pPr_to_copy.attrib)
                            
                            # Copy all child elements (like <a:buFont>, <a:buChar>, etc.)
                            for child in pPr_to_copy:
                                new_pPr.append(deepcopy(child))
                        
                        
                        # Add the text with the original font style
                        run = new_p.add_run()
                        run.text = item_text
                        if source_font:
                            _transfer_font_properties(source_font, run.font)
                            
                else: # Handle empty list or invalid data type
                    # Set the original paragraph text to "None"
                    p.clear()
                    run = p.add_run()
                    run.text = "None"
                    if source_font:
                        _transfer_font_properties(source_font, run.font)

                # --- Text Frame Properties ---
                tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                tf.word_wrap = True

                found_list_in_shape = True # Mark that we handled a list in this shape

            if found_list_in_shape:
                continue # Skip standard text replacement for this shape

            # --- Text Replacement Logic (preserving formatting) ---
                        # --- Text Replacement Logic (preserving formatting) ---
            for para in shape.text_frame.paragraphs:
                # Optimization: Skip paragraphs that don't contain any placeholders
                if '{{' not in para.text:
                    continue

                # --- Attempt 1: Run-by-Run Replacement (Preserves Formatting) ---
                was_run_replacement_made = False
                for run in para.runs:
                    if '{{' not in run.text:
                        continue
                    
                    matches = text_pattern.findall(run.text)
                    if not matches:
                        continue

                    modified_text = run.text
                    for ph_name in matches:
                        replacement_value = str(data.get(ph_name, ""))
                        
                        placeholder_tag_text = f"{{{{text:{ph_name}}}}}"
                        placeholder_tag_choice = f"{{{{choice:{ph_name}}}}}"
                        placeholder_tag_simple = f"{{{{{ph_name}}}}}"
                        
                        modified_text = modified_text.replace(placeholder_tag_text, replacement_value)
                        modified_text = modified_text.replace(placeholder_tag_choice, replacement_value)
                        modified_text = modified_text.replace(placeholder_tag_simple, replacement_value)
                    
                    if modified_text != run.text:
                        run.text = modified_text
                        was_run_replacement_made = True

                # --- Attempt 2: Fallback for Split-Run Placeholders ---
                # If no runs were replaced, but the paragraph *still* has a
                # placeholder, it must be split across runs.
                if not was_run_replacement_made and '{{' in para.text:
                    
                    # We must use the "whole paragraph" method.
                    full_text_from_runs = "".join(run.text for run in para.runs)
                    matches = text_pattern.findall(full_text_from_runs)

                    if not matches:
                        continue # Should be rare, but a safe check

                    source_font = para.runs[0].font if para.runs else None
                    modified_full_text = full_text_from_runs
                    
                    for ph_name in matches:
                        replacement_value = str(data.get(ph_name, ""))
                        
                        placeholder_tag_text = f"{{{{text:{ph_name}}}}}"
                        placeholder_tag_choice = f"{{{{choice:{ph_name}}}}}"
                        placeholder_tag_simple = f"{{{{{ph_name}}}}}"
                        
                        modified_full_text = modified_full_text.replace(placeholder_tag_text, replacement_value)
                        modified_full_text = modified_full_text.replace(placeholder_tag_choice, replacement_value)
                        modified_full_text = modified_full_text.replace(placeholder_tag_simple, replacement_value)
                    
                    para.clear()
                    new_run = para.add_run()
                    new_run.text = modified_full_text
                    
                    if source_font:
                        _transfer_font_properties(source_font, new_run.font)
                        
        # After iterating all shapes, delete the placeholder shapes
        for shape in shapes_to_delete:
            sp_element = shape.element
            sp_element.getparent().remove(sp_element)

    # Save the final presentation to a new in-memory stream
    output_stream = BytesIO()
    ppt.save(output_stream)
    output_stream.seek(0)
    
    return output_stream
